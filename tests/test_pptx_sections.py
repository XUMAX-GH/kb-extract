"""SP-2 / v0.2.0: PPTX native Sections form a root → section → slide tree.

Native PowerPoint Sections are stored in `ppt/presentation.xml` as
`<p14:sectionLst>` and have no public python-pptx API; we inject the XML
directly to build a deterministic fixture.
"""

from __future__ import annotations

from pathlib import Path

import pptx
import pytest
from lxml import etree

from kb_extract.adapters.pptx import PptxAdapter


def _make_sectioned_pptx(path: Path, sections: list[tuple[str, int]]) -> None:
    """Create a .pptx whose slides are grouped into the given named sections.

    sections: [(section_name, slide_count_in_section), ...]
    """
    prs = pptx.Presentation()
    blank_layout = prs.slide_layouts[6]
    # Add the total slides
    total = sum(c for _, c in sections)
    for _ in range(total):
        prs.slides.add_slide(blank_layout)
    slide_ids = [s.slide_id for s in prs.slides]

    # Build the <p:extLst><p:ext><p14:sectionLst>...</p:extLst> subtree
    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    ns_p14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    nsmap = {"p": ns_p, "p14": ns_p14}
    extLst = etree.SubElement(prs.element, etree.QName(ns_p, "extLst"), nsmap=nsmap)
    ext = etree.SubElement(
        extLst, etree.QName(ns_p, "ext"),
        attrib={"uri": "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"},
    )
    section_list = etree.SubElement(
        ext, etree.QName(ns_p14, "sectionLst"), nsmap={"p14": ns_p14},
    )
    cursor = 0
    for sec_idx, (sec_name, count) in enumerate(sections, start=1):
        section = etree.SubElement(
            section_list, etree.QName(ns_p14, "section"),
            attrib={
                "name": sec_name,
                "id": f"{{00000000-0000-0000-0000-{sec_idx:012d}}}",
            },
        )
        sld_id_lst = etree.SubElement(section, etree.QName(ns_p14, "sldIdLst"))
        for _ in range(count):
            etree.SubElement(
                sld_id_lst, etree.QName(ns_p14, "sldId"),
                attrib={"id": str(slide_ids[cursor])},
            )
            cursor += 1

    prs.save(str(path))


@pytest.mark.disable_socket
def test_pptx_adapter_builds_section_tree_when_sections_present(tmp_path):
    src = tmp_path / "deck.pptx"
    _make_sectioned_pptx(src, [("Intro", 1), ("Body", 2)])
    out = tmp_path / "out.tmp"
    out.mkdir()

    result = PptxAdapter().extract(src, out)

    # outline_source bumps to pptx_section when sections are detected
    assert result.meta.outline_source == "pptx_section"
    assert result.meta.outline_confidence == "high"

    # Root has two section children (Intro, Body)
    sections = result.index.children
    assert len(sections) == 2
    assert sections[0].title == "Intro"
    assert sections[1].title == "Body"
    # Section is level=1, slides under it are level=2
    assert sections[0].level == 1
    assert sections[1].level == 1
    # Intro contains 1 slide; Body contains 2
    assert len(sections[0].children) == 1
    assert len(sections[1].children) == 2
    assert all(c.level == 2 for c in sections[0].children)
    assert all(c.level == 2 for c in sections[1].children)
    # Slides under Body are pages 2 and 3
    body_pages = [c.page_start for c in sections[1].children]
    assert body_pages == [2, 3]


@pytest.mark.disable_socket
def test_pptx_adapter_falls_back_to_flat_when_no_sections(tmp_path):
    """A .pptx without sectionLst keeps the v0.1 flat behavior."""
    prs = pptx.Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.slides.add_slide(prs.slide_layouts[6])
    src = tmp_path / "deck.pptx"
    prs.save(str(src))

    out = tmp_path / "out.tmp"
    out.mkdir()
    result = PptxAdapter().extract(src, out)

    assert result.meta.outline_source == "heading_style"
    # Flat children, each level=1
    assert all(c.level == 1 for c in result.index.children)
    assert len(result.index.children) == 2
