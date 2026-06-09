"""SP-2 / v0.2.0: every adapter writes a valid outline_confidence on meta."""

from __future__ import annotations

import json
from pathlib import Path

import openpyxl
import pptx
import pytest


@pytest.mark.disable_socket
def test_xlsx_meta_has_outline_confidence_high(tmp_path):
    from kb_extract.adapters.xlsx import XlsxAdapter

    wb = openpyxl.Workbook()
    wb.active.title = "Only"
    wb["Only"]["A1"] = "x"
    src = tmp_path / "x.xlsx"
    wb.save(str(src))
    out = tmp_path / "o.tmp"
    out.mkdir()
    r = XlsxAdapter().extract(src, out)
    assert r.meta.outline_confidence == "high"


@pytest.mark.disable_socket
def test_pptx_meta_has_outline_confidence_high(tmp_path):
    from kb_extract.adapters.pptx import PptxAdapter

    prs = pptx.Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    src = tmp_path / "d.pptx"
    prs.save(str(src))
    out = tmp_path / "o.tmp"
    out.mkdir()
    r = PptxAdapter().extract(src, out)
    assert r.meta.outline_confidence == "high"


@pytest.mark.disable_socket
def test_meta_json_serializes_outline_confidence(tmp_path):
    """End-to-end through the orchestrator: meta.json on disk contains the field."""
    from kb_extract.adapters._noop import NoopAdapter
    from kb_extract.adapters.base import Registry
    from kb_extract.orchestrator import run

    project = tmp_path / "P"
    project.mkdir()
    (project / "a.noop").write_bytes(b"hi")

    reg = Registry()
    reg.register(NoopAdapter())
    run(project, registry=reg)

    meta_files = list((project / "kb").glob("**/meta.json"))
    assert meta_files, "no meta.json emitted"
    data = json.loads(Path(meta_files[0]).read_text(encoding="utf-8"))
    assert "outline_confidence" in data
    assert data["outline_confidence"] in {"high", "medium", "low"}
