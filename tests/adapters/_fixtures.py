"""Helpers that generate tiny synthetic fixtures on the fly.

Avoids committing binary blobs; matches the "no Microsoft confidential docs"
rule in tests/fixtures/SOURCES.md.
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image as PILImage


def make_png(path: Path, w: int = 4, h: int = 3, color=(255, 0, 0)) -> Path:
    img = PILImage.new("RGB", (w, h), color)
    img.save(path, format="PNG")
    return path


def make_docx(path: Path) -> Path:
    from docx import Document
    doc = Document()
    doc.add_heading("Chapter 1", level=1)
    doc.add_paragraph("First paragraph in chapter 1.")
    doc.add_heading("Section 1.1", level=2)
    doc.add_paragraph("Body of section 1.1.")
    t = doc.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "col A"
    t.cell(0, 1).text = "col B"
    t.cell(1, 0).text = "1"
    t.cell(1, 1).text = "2"
    doc.save(str(path))
    return path


def make_xlsx(path: Path) -> Path:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Summary"
    ws["A1"] = "metric"
    ws["B1"] = "value"
    ws["A2"] = "count"
    ws["B2"] = 42
    ws["A3"] = "ratio"
    ws["B3"] = 0.5
    ws2 = wb.create_sheet("Details")
    ws2["A1"] = "id"
    ws2["B1"] = "name"
    ws2["A2"] = 1
    ws2["B2"] = "alpha"
    wb.save(str(path))
    return path


def make_pptx(path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # title slide
    slide1.shapes.title.text = "First Slide"
    if slide1.placeholders[1].has_text_frame:
        slide1.placeholders[1].text = "Subtitle text"
    slide1.notes_slide.notes_text_frame.text = "presenter note one"

    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide2.shapes.title.text = "Second Slide"
    tx_box = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    tx_box.text_frame.text = "Some bullet body text on slide 2"
    prs.save(str(path))
    return path
