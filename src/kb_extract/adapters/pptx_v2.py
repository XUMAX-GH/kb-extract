"""PPTX v2 adapter — MinerU-inspired (v0.8.0 parser v2, spec sec.4.2).

Scope of this commit:
- Image extraction: ``MSO_SHAPE_TYPE.PICTURE`` shapes -> ``assets/`` via
  :mod:`._image_utils`, with ``![](assets/slide_N_img_M.ext)`` markdown
  links inserted in document order.
- Table cells -> HTML ``<table>`` via :mod:`._table_utils`, preserving
  merged cells via PPTX ``cell.span_height`` / ``cell.span_width``.
- GroupShape recursion: text and images inside groups are walked.
- Speaker notes: kept (already worked in legacy) but switched to the
  ``> **Note:**`` markdown marker per spec.

Out of scope: SmartArt content extraction (spec explicit exclusion).

Section detection, anchor scheme, page count, and hardness guarantees
match the legacy ``PptxAdapter``.
"""
from __future__ import annotations

import hashlib
from pathlib import Path

import pptx as _pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..contracts import ExtractionResult, SectionNode, TableRef
from ._common import make_meta
from ._image_utils import save_image
from ._table_utils import CellInfo, cells_to_html

_PPTX_VERSION = getattr(_pptx, "__version__", "unknown")

_NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
}


def _detect_sections(prs) -> list[tuple[str, list[int]]] | None:
    """Same logic as legacy PptxAdapter: read p14:sectionLst from raw XML."""
    root = prs.element
    section_list = root.find(".//p14:sectionLst", _NS)
    if section_list is None:
        return None
    id_to_idx: dict[int, int] = {}
    for idx, slide in enumerate(prs.slides, start=1):
        try:
            id_to_idx[slide.slide_id] = idx
        except Exception:
            continue
    result: list[tuple[str, list[int]]] = []
    for section in section_list.findall("p14:section", _NS):
        name = section.get("name") or "(Unnamed Section)"
        indices: list[int] = []
        for sid_elem in section.findall("p14:sldIdLst/p14:sldId", _NS):
            raw = sid_elem.get("id")
            if raw is None:
                continue
            try:
                sid = int(raw)
            except ValueError:
                continue
            if sid in id_to_idx:
                indices.append(id_to_idx[sid])
        if indices:
            result.append((name, sorted(indices)))
    return result if result else None


def _pptx_table_to_grid(table) -> list[list[CellInfo]]:
    """Convert a python-pptx Table into a phantom-free grid.

    python-pptx exposes ``cell.is_spanned`` (True for cells covered by a
    merge to their left/above), ``cell.is_merge_origin`` (True for the
    top-left), ``cell.span_width`` (colspan) and ``cell.span_height``
    (rowspan).
    """
    out: list[list[CellInfo]] = []
    for row in table.rows:
        out_row: list[CellInfo] = []
        for cell in row.cells:
            # Drop continuation cells; only origin cells emit a CellInfo
            if cell.is_spanned and not cell.is_merge_origin:
                continue
            text = (cell.text or "").strip()
            colspan = max(1, getattr(cell, "span_width", 1) or 1)
            rowspan = max(1, getattr(cell, "span_height", 1) or 1)
            out_row.append(CellInfo(text=text, colspan=colspan, rowspan=rowspan))
        out.append(out_row)
    return out


def _walk_shapes(shapes):
    """Yield shapes recursively, descending into GroupShapes."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes)
        else:
            yield shape


def _render_slide_content(
    slide,
    slide_idx: int,
    out_dir: Path,
    image_counter: list[int],
    table_counter: list[int],
    tables_out: list[TableRef],
) -> list[str]:
    """Render all non-title shapes on a slide to markdown lines."""
    title_shape = slide.shapes.title
    chunks: list[str] = []

    img_idx_in_slide = 0
    for shape in _walk_shapes(slide.shapes):
        if shape is title_shape:
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
            except Exception:
                continue
            img_idx_in_slide += 1
            image_counter[0] += 1
            rel = save_image(
                blob, out_dir,
                prefix=f"slide_{slide_idx}_img",
                index=img_idx_in_slide,
            )
            if rel is not None:
                chunks.append(f"![]({rel})")
            continue

        if shape.has_table:
            table_counter[0] += 1
            anchor = f"tbl-{table_counter[0]:04d}"
            grid = _pptx_table_to_grid(shape.table)
            rows_json: tuple[tuple[str, ...], ...] = tuple(
                tuple(c.text for c in row) for row in grid
            )
            chunks.append(f'<a id="{anchor}"></a>')
            chunks.append(cells_to_html(grid))
            tables_out.append(TableRef(
                anchor=anchor, page=slide_idx, rows_json=rows_json,
                rendered_asset=None,
            ))
            continue

        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt:
                chunks.append(txt)

    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            for line in notes.splitlines():
                chunks.append(f"> **Note:** {line}")

    return chunks


class PptxV2Adapter:
    name = "pptx_v2"
    version = "0.8.0"
    extensions = (".pptx",)

    def extract(self, src: Path, out_dir_tmp: Path) -> ExtractionResult:
        prs = _pptx.Presentation(str(src))
        sha = hashlib.sha256(src.read_bytes()).hexdigest()
        md_lines: list[str] = [
            f"<!-- generated by kb-extract; do not edit; source_sha256: {sha} -->",
            "",
        ]

        image_counter = [0]  # boxed for closure mutation
        table_counter = [0]
        tables_out: list[TableRef] = []

        slide_records: list[tuple[str, str, list[str]]] = []
        for i, slide in enumerate(prs.slides, start=1):
            title = ""
            if slide.shapes.title is not None:
                title = (slide.shapes.title.text or "").strip()
            if not title:
                title = f"Slide {i}"
            anchor = f"sec-{i:04d}"
            chunks = _render_slide_content(
                slide, i, out_dir_tmp,
                image_counter, table_counter, tables_out,
            )
            slide_records.append((title, anchor, chunks))

        sections = _detect_sections(prs)
        children: list[SectionNode] = []
        if sections:
            outline_source = "pptx_section"
            for sec_counter, (sec_name, slide_indices) in enumerate(sections, start=1):
                sec_anchor = f"sec-grp-{sec_counter:04d}"
                md_lines.append(f'<a id="{sec_anchor}"></a>')
                md_lines.append(f"# {sec_name}")
                md_lines.append("")
                grandchildren: list[SectionNode] = []
                for slide_idx in slide_indices:
                    title, anchor, chunks = slide_records[slide_idx - 1]
                    md_lines.append(f'<a id="{anchor}"></a>')
                    md_lines.append(f"## {title}")
                    md_lines.append("")
                    for c in chunks:
                        md_lines.append(c)
                        md_lines.append("")
                    grandchildren.append(SectionNode(
                        node_id=f"{sec_counter:04d}-{slide_idx:04d}",
                        title=title, level=2,
                        page_start=slide_idx, page_end=slide_idx,
                        anchor=anchor, language="und",
                    ))
                children.append(SectionNode(
                    node_id=f"{sec_counter:04d}",
                    title=sec_name, level=1,
                    page_start=slide_indices[0], page_end=slide_indices[-1],
                    anchor=sec_anchor, language="und",
                    children=tuple(grandchildren),
                ))
        else:
            outline_source = "heading_style"
            for i, (title, anchor, chunks) in enumerate(slide_records, start=1):
                md_lines.append(f'<a id="{anchor}"></a>')
                md_lines.append(f"# {title}")
                md_lines.append("")
                for c in chunks:
                    md_lines.append(c)
                    md_lines.append("")
                children.append(SectionNode(
                    node_id=f"{i:04d}", title=title, level=1,
                    page_start=i, page_end=i, anchor=anchor, language="und",
                ))

        total = max(len(slide_records), 1)
        root = SectionNode(
            node_id="0000", title=src.stem, level=0,
            page_start=1, page_end=total,
            anchor="", language="und", children=tuple(children),
        )
        markdown = "\n".join(md_lines) + "\n"
        meta = make_meta(
            src=src,
            adapter_name=self.name,
            adapter_version=self.version,
            tool_versions={"python-pptx": _PPTX_VERSION},
            outline_source=outline_source,  # type: ignore[arg-type]
            outline_confidence="high",
        )
        return ExtractionResult(
            markdown=markdown, index=root,
            tables=tuple(tables_out), assets=(), meta=meta,
        )
